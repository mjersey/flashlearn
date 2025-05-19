import os
import fitz  # PyMuPDF
import pptx
from pptx import Presentation
import requests
import tkinter as tk
from tkinter import filedialog
import re

# Hugging Face API Setup
HUGGINGFACE_API_KEY = "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"
API_URL = "https://api-inference.huggingface.co/models/valhalla/t5-base-qg-hl"
HEADERS = {"Authorization": f"Bearer {HUGGINGFACE_API_KEY}"}


# Function to extract text from PDF
def extract_text_from_pdf(pdf_path):
    doc = fitz.open(pdf_path)
    text = "\n".join([page.get_text("text") for page in doc])
    return text.encode('utf-8', 'ignore').decode('utf-8')


# Function to extract text from PPTX
def extract_text_from_pptx(pptx_path):
    prs = pptx.Presentation(pptx_path)
    text = "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
    return text.encode('utf-8').decode('utf-8')


# Function to split text into smaller chunks
def split_text(text, chunk_size=500):
    words = text.split()
    return [" ".join(words[i:i + chunk_size]) for i in range(0, len(words), chunk_size)]


# Function to find the best matching answer from the text
def find_answer_in_text(question, text):
    sentences = re.split(r'(?<=[.!?]) +', text)  # Split text into sentences
    question_keywords = set(question.lower().split())

    best_match = None
    highest_overlap = 0

    for sentence in sentences:
        sentence_words = set(sentence.lower().split())
        common_words = question_keywords.intersection(sentence_words)

        if len(common_words) > highest_overlap:
            highest_overlap = len(common_words)
            best_match = sentence

    return best_match if best_match else "Answer not found in text."


# Function to generate flashcards (Q&A)
def generate_flashcards(text):
    chunks = split_text(text, chunk_size=150)
    flashcards = []

    for chunk in chunks:
        # Generate a question using Hugging Face API
        payload = {"inputs": chunk}
        response = requests.post(API_URL, headers=HEADERS, json=payload)

        if response.status_code != 200:
            continue

        try:
            response_data = response.json()
            if isinstance(response_data, list) and len(response_data) > 0:
                for item in response_data:
                    question = item.get("generated_text", "No question generated.")

                    # Find the best possible answer
                    correct_answer = find_answer_in_text(question, text)

                    # Store question and answer in a dictionary
                    flashcards.append({"question": question, "answer": correct_answer})
        except Exception:
            continue

    return flashcards


# Function to handle file upload
def upload_file():
    global flashcards, current_index

    file_path = filedialog.askopenfilename(filetypes=[("PDF & PPT Files", "*.pdf;*.pptx")])
    if file_path:
        if file_path.endswith(".pdf"):
            text = extract_text_from_pdf(file_path)
        elif file_path.endswith(".pptx"):
            text = extract_text_from_pptx(file_path)
        else:
            flashcard_label.config(text="Unsupported file format.")
            return

        if not text.strip():
            flashcard_label.config(text="Error: The file contains no readable text.")
            return

        flashcards = generate_flashcards(text)
        current_index = 0
        show_flashcard()


# Function to show flashcard (hide answer initially)
def show_flashcard():
    global current_index, answer_hidden
    if flashcards:
        flashcard_label.config(text=f"Q: {flashcards[current_index]['question']}")
        answer_label.config(text="")  # Hide answer
        answer_hidden = True
    else:
        flashcard_label.config(text="No flashcards generated.")
        answer_label.config(text="")


# Function to reveal the answer
def reveal_answer():
    global answer_hidden
    if flashcards and answer_hidden:
        answer_label.config(text=f"A: {flashcards[current_index]['answer']}")
        answer_hidden = False


# Function to move to the next flashcard (hide answer again)
def next_flashcard():
    global current_index, answer_hidden
    if flashcards:
        current_index = (current_index + 1) % len(flashcards)
        show_flashcard()


# Tkinter UI Setup
root = tk.Tk()
root.title("AI Flashcard Generator")
root.geometry("600x400")
root.configure(bg="#f4f4f4")

# Upload Button
upload_button = tk.Button(root, text="Upload PDF/PPT", command=upload_file, font=("Arial", 12), bg="#4CAF50", fg="white", padx=10, pady=5)
upload_button.pack(pady=10)

# Flashcard Display Area
flashcard_label = tk.Label(root, text="Upload a file to generate flashcards.", wraplength=500, font=("Arial", 14), bg="white", relief="solid", bd=2, padx=20, pady=20)
flashcard_label.pack(pady=10, padx=20, fill="both", expand=True)

# Answer Display Area (Hidden Initially)
answer_label = tk.Label(root, text="", wraplength=500, font=("Arial", 14), fg="blue", bg="white", padx=10, pady=10)
answer_label.pack()

# Show Answer Button
show_answer_button = tk.Button(root, text="Show Answer", command=reveal_answer, font=("Arial", 12), bg="#FFA500", fg="white", padx=10, pady=5)
show_answer_button.pack(pady=5)

# Next Flashcard Button
next_button = tk.Button(root, text="Next", command=next_flashcard, font=("Arial", 12), bg="#008CBA", fg="white", padx=10, pady=5)
next_button.pack(pady=10)

# Run Tkinter App
flashcards = []
current_index = 0
answer_hidden = True
root.mainloop()
