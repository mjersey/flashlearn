import subprocess
from PIL import Image, ImageTk, ImageDraw
import os
import textwrap
import tkinter as tk
from tkinter import filedialog, messagebox, scrolledtext, ttk
import fitz
import pptx
import json
import time
import random
from settings import get_current_settings, apply_theme, apply_font_size, THEMES, FONT_SIZES

flashcards = []
current_index = 0
answer_hidden = True
decks = []
extracted_text = ""

DECKS_FOLDER = "user_decks"
os.makedirs(DECKS_FOLDER, exist_ok=True)

# Constants
PROGRESS_FOLDER = "user_progress"
SETTINGS_FOLDER = "user_settings"
os.makedirs(PROGRESS_FOLDER, exist_ok=True)
os.makedirs(SETTINGS_FOLDER, exist_ok=True)


def center_window(window, width, height):
    """Center a window on the screen"""
    # Get screen dimensions
    screen_width = window.winfo_screenwidth()
    screen_height = window.winfo_screenheight()

    # Calculate position
    x = (screen_width - width) // 2
    y = (screen_height - height) // 2

    # Set window position
    window.geometry(f"{width}x{height}+{x}+{y}")


def get_user_decks_file(username):
    """Get the path to the user's decks file"""
    return os.path.join(DECKS_FOLDER, f"{username}_decks.json")


def load_decks():
    """Load saved decks from file for the current user"""
    global decks
    username = get_logged_in_user()[0]
    decks_file = get_user_decks_file(username)

    if os.path.exists(decks_file):
        try:
            with open(decks_file, 'r') as f:
                decks = json.load(f)
        except Exception as e:
            print(f"Error loading decks: {e}")
            decks = []
    else:
        decks = []
    return decks


def save_decks():
    """Save decks to file for the current user"""
    username = get_logged_in_user()[0]
    decks_file = get_user_decks_file(username)

    try:
        with open(decks_file, 'w') as f:
            json.dump(decks, f)
    except Exception as e:
        print(f"Error saving decks: {e}")
        messagebox.showerror("Error", f"Failed to save decks: {str(e)}")


def open_page(page):
    root.destroy()
    subprocess.run(["python", f"{page}.py"])


def sign_out():
    if os.path.exists("current_user.txt"):
        os.remove("current_user.txt")
    root.destroy()
    import sign_in


PROFILE_PICTURE_FOLDER = "profile_pictures"
os.makedirs(PROFILE_PICTURE_FOLDER, exist_ok=True)


def get_logged_in_user():
    try:
        with open("current_user.txt", "r") as user_file:
            data = user_file.read().strip().split(",")
            return (data + [None])[:4]  # Ensure four elements (username, email, password, profile_pic)
    except FileNotFoundError:
        return "Guest", "No Email", "No Password", None


def get_user_progress(username):
    progress_file = os.path.join(PROGRESS_FOLDER, f"{username}_progress.json")
    if os.path.exists(progress_file):
        try:
            with open(progress_file, 'r') as f:
                return json.load(f)
        except Exception as e:
            print(f"Error loading progress: {e}")
    return {}


def save_user_progress(username, progress_data):
    progress_file = os.path.join(PROGRESS_FOLDER, f"{username}_progress.json")
    try:
        with open(progress_file, 'w') as f:
            json.dump(progress_data, f)
        return True
    except Exception as e:
        print(f"Error saving progress: {e}")
        return False


def update_card_progress(deck_title, card_index, is_correct=True):
    """Update progress for a specific card in a deck"""
    username = get_logged_in_user()[0]
    progress_data = get_user_progress(username)

    # Initialize progress data for this deck if it doesn't exist
    if deck_title not in progress_data:
        progress_data[deck_title] = {
            'view_count': 0,
            'cards_viewed': 0,
            'correct_answers': 0,
            'last_viewed': '',
            'card_status': {}  # Track status of individual cards
        }

    # Update view count and last viewed timestamp
    progress_data[deck_title]['view_count'] += 1
    progress_data[deck_title]['last_viewed'] = time.strftime("%Y-%m-%d %H:%M:%S")

    # Get the current card status dictionary
    card_status = progress_data[deck_title].get('card_status', {})

    # Convert card_index to string for JSON compatibility
    card_index_str = str(card_index)

    # If this is the first time viewing this card, increment cards_viewed
    if card_index_str not in card_status:
        progress_data[deck_title]['cards_viewed'] += 1

    # Update the card's status based on the current answer
    previous_status = card_status.get(card_index_str)

    # If the card was previously correct and is now wrong, decrement correct_answers
    if previous_status == 'correct' and not is_correct:
        progress_data[deck_title]['correct_answers'] = max(0, progress_data[deck_title]['correct_answers'] - 1)

    # If the card was previously wrong and is now correct, increment correct_answers
    elif previous_status == 'wrong' and is_correct:
        progress_data[deck_title]['correct_answers'] += 1

    # If this is the first time answering this card and it's correct, increment correct_answers
    elif previous_status is None and is_correct:
        progress_data[deck_title]['correct_answers'] += 1

    # Update the card's status
    card_status[card_index_str] = 'correct' if is_correct else 'wrong'
    progress_data[deck_title]['card_status'] = card_status

    # Save updated progress
    save_user_progress(username, progress_data)

    # Update the dashboard progress display
    if 'progress_canvas' in globals():
        update_progress_display()


def extract_text_from_pdf(pdf_path):
    try:
        doc = fitz.open(pdf_path)
        text = "\n".join([page.get_text("text") for page in doc])
        return text
    except Exception as e:
        messagebox.showerror("Error", f"Failed to extract text from PDF: {str(e)}")
        return ""


def extract_text_from_pptx(pptx_path):
    try:
        prs = pptx.Presentation(pptx_path)
        text = "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
        return text
    except Exception as e:
        messagebox.showerror("Error", f"Failed to extract text from PPTX: {str(e)}")
        return ""


def upload_file(text_preview=None):
    global extracted_text
    file_path = filedialog.askopenfilename(filetypes=[("PDF & PPT Files", "*.pdf;*.pptx")])

    if file_path:
        if file_path.endswith(".pdf"):
            extracted_text = extract_text_from_pdf(file_path)
        else:
            extracted_text = extract_text_from_pptx(file_path)

        if not extracted_text.strip():
            messagebox.showerror("Error", "The file contains no readable text.")
            return

        # Show preview of extracted text if text_preview widget is provided
        if text_preview:
            text_preview.delete(1.0, tk.END)
            text_preview.insert(tk.END, extracted_text[:1000] + "..." if len(extracted_text) > 1000 else extracted_text)

        messagebox.showinfo("Success",
                            "File uploaded successfully. You can use this text as reference for creating cards.")


def show_flashcard():
    global current_index, answer_hidden
    if flashcards:
        flashcard_label.config(text=f"Q: {flashcards[current_index]['question']}")
        answer_label.config(text="")
        answer_hidden = True
    else:
        flashcard_label.config(text="No flashcards generated.")
        answer_label.config(text="")


def reveal_answer():
    global answer_hidden
    if flashcards and answer_hidden:
        answer_label.config(text=f"A: {flashcards[current_index]['answer']}")
        answer_hidden = False


def next_flashcard():
    global current_index
    if flashcards:
        current_index = (current_index + 1) % len(flashcards)
        show_flashcard()


def delete_deck(deck_index):
    """Delete a deck from the list and update the display"""
    if 0 <= deck_index < len(decks):
        # Ask for confirmation
        if messagebox.askyesno("Confirm Delete",
                               f"Are you sure you want to delete the deck '{decks[deck_index]['title']}'?"):
            # Remove the deck
            del decks[deck_index]
            # Save the updated decks list
            save_decks()
            # Update the display
            update_deck_display()
            # Update progress display
            update_progress_display()
            messagebox.showinfo("Success", "Deck deleted successfully.")


def open_card_window(edit_deck=None, edit_index=None):
    """Open card window for creating or editing a deck"""
    global extracted_text
    extracted_text = ""  # Reset extracted text

    # Create the window
    card_window = tk.Toplevel(root)
    card_window.title("Add New Card" if not edit_deck else "Edit Deck")
    card_window.geometry("550x650")
    card_window.configure(bg="white")

    # Center the window
    center_window(card_window, 550, 650)

    # Make window modal (user must interact with this window first)
    card_window.grab_set()
    card_window.transient(root)

    # Main content frame with padding
    content_frame = tk.Frame(card_window, bg="white", padx=20, pady=20)
    content_frame.pack(fill="both", expand=True)

    # DECK TITLE section
    title_label = tk.Label(content_frame, text="DECK TITLE:", font=("Arial", 11, "bold"), bg="white", anchor="w")
    title_label.pack(fill="x", anchor="w")

    title_entry = tk.Entry(content_frame, font=("Arial", 11), bd=1, relief="solid")
    title_entry.pack(fill="x", pady=(5, 15))

    # If editing, populate the title
    if edit_deck:
        title_entry.insert(0, edit_deck["title"])

    # REFERENCE TEXT section
    ref_frame = tk.Frame(content_frame, bg="white")
    ref_frame.pack(fill="x", pady=(5, 15))

    ref_label = tk.Label(ref_frame, text="REFERENCE TEXT (OPTIONAL):", font=("Arial", 11, "bold"), bg="white",
                         anchor="w")
    ref_label.pack(side="left", fill="x", anchor="w")

    import_btn = tk.Button(
        ref_frame,
        text="Import Reference File",
        font=("Arial", 10),
        bg="#003366",
        fg="white",
        padx=10,
        pady=2,
        relief="flat",
        command=lambda: upload_file(ref_text)
    )
    import_btn.pack(side="right", padx=(10, 0))

    ref_text_frame = tk.Frame(content_frame, bg="white")
    ref_text_frame.pack(fill="x", pady=(0, 15))

    ref_text = tk.Text(ref_text_frame, height=8, font=("Arial", 11), bd=1, relief="solid")
    ref_text.pack(side="left", fill="both", expand=True)

    ref_scrollbar = tk.Scrollbar(ref_text_frame, command=ref_text.yview)
    ref_scrollbar.pack(side="right", fill="y")
    ref_text.config(yscrollcommand=ref_scrollbar.set)

    # CARDS section
    cards_frame = tk.Frame(content_frame, bg="white")
    cards_frame.pack(fill="x", pady=(5, 15))

    cards_label = tk.Label(cards_frame, text="CARDS:", font=("Arial", 11, "bold"), bg="white", anchor="w")
    cards_label.pack(side="left", fill="x", anchor="w")

    add_card_btn = tk.Button(
        cards_frame,
        text="+ Add Another Card",
        font=("Arial", 10),
        bg="#003366",
        fg="white",
        padx=10,
        pady=2,
        relief="flat"
    )
    add_card_btn.pack(side="right", padx=(10, 0))

    # Frame to hold all card entries with navigation
    cards_container_frame = tk.Frame(content_frame, bg="white")
    cards_container_frame.pack(fill="x", pady=(0, 15))

    # Navigation frame
    nav_frame = tk.Frame(cards_container_frame, bg="white")
    nav_frame.pack(fill="x", pady=(5, 0))

    # Card display frame
    card_display_frame = tk.Frame(cards_container_frame, bg="#f0f0f0", bd=1, relief="solid")
    card_display_frame.pack(fill="x", pady=(5, 10))

    # List to store card entry widgets
    card_entries = []
    current_card_index = [0]  # Using a list to make it mutable in nested functions

    def create_card_entry(question="", answer=""):
        """Create a card entry frame with question and answer fields"""
        card_frame = tk.Frame(card_display_frame, bg="#f0f0f0")

        # Question label and entry
        question_label = tk.Label(card_frame, text="Question:", font=("Arial", 10), bg="#f0f0f0", anchor="w")
        question_label.pack(fill="x", anchor="w", padx=10, pady=(10, 0))

        question_entry = tk.Entry(card_frame, font=("Arial", 10), width=50)
        question_entry.pack(fill="x", padx=10, pady=(2, 5))
        if question:
            question_entry.insert(0, question)

        # Answer label and entry
        answer_label = tk.Label(card_frame, text="Answer:", font=("Arial", 10), bg="#f0f0f0", anchor="w")
        answer_label.pack(fill="x", anchor="w", padx=10)

        answer_entry = tk.Entry(card_frame, font=("Arial", 10), width=50)
        answer_entry.pack(fill="x", padx=10, pady=(2, 5))
        if answer:
            answer_entry.insert(0, answer)

        # Remove button (aligned to the right)
        remove_btn_frame = tk.Frame(card_frame, bg="#f0f0f0")
        remove_btn_frame.pack(fill="x", padx=10, pady=(5, 10))

        remove_btn = tk.Button(
            remove_btn_frame,
            text="Remove",
            font=("Arial", 9),
            bg="#FF0000",  # Changed to red
            fg="white",
            padx=5,
            pady=1,
            relief="flat"
        )
        remove_btn.pack(side="right")

        return card_frame, question_entry, answer_entry, remove_btn

    def add_card_entry(question="", answer=""):
        """Add a new card entry to the list"""
        card_frame, question_entry, answer_entry, remove_btn = create_card_entry(question, answer)
        card_entry = {
            "frame": card_frame,
            "question": question_entry,
            "answer": answer_entry
        }
        card_entries.append(card_entry)

        # Configure remove button to remove this specific card
        remove_btn.config(command=lambda: remove_card_entry(card_entry))

        # Show the latest card
        current_card_index[0] = len(card_entries) - 1
        show_current_card()

        # Update navigation buttons
        update_nav_buttons()

    def remove_card_entry(card_entry):
        """Remove a card entry from the list"""
        if len(card_entries) <= 1:
            messagebox.showinfo("Info", "You must have at least one card.")
            return

        # Remove from the list
        idx = card_entries.index(card_entry)
        card_entries.remove(card_entry)

        # Adjust current index if needed
        if current_card_index[0] >= len(card_entries):
            current_card_index[0] = len(card_entries) - 1

        # Show the current card
        show_current_card()

        # Update navigation buttons
        update_nav_buttons()

    def show_current_card():
        """Show the current card and hide others"""
        # Hide all card frames
        for entry in card_entries:
            entry["frame"].pack_forget()

        # Show current card frame
        if card_entries:
            card_entries[current_card_index[0]]["frame"].pack(fill="x")

            # Update card counter
            card_counter_label.config(text=f"Card {current_card_index[0] + 1} of {len(card_entries)}")

    def next_card():
        """Show the next card"""
        if current_card_index[0] < len(card_entries) - 1:
            current_card_index[0] += 1
            show_current_card()
            update_nav_buttons()

    def prev_card():
        """Show the previous card"""
        if current_card_index[0] > 0:
            current_card_index[0] -= 1
            show_current_card()
            update_nav_buttons()

    def update_nav_buttons():
        """Update the state of navigation buttons"""
        # Disable prev button if at first card
        if current_card_index[0] == 0:
            prev_btn.config(state="disabled")
        else:
            prev_btn.config(state="normal")

        # Disable next button if at last card
        if current_card_index[0] >= len(card_entries) - 1:
            next_btn.config(state="disabled")
        else:
            next_btn.config(state="normal")

    # Navigation buttons and card counter
    prev_btn = tk.Button(
        nav_frame,
        text="◀ Previous",
        font=("Arial", 9),
        bg="#999999",  # Changed to gray
        fg="white",
        padx=5,
        pady=1,
        relief="flat",
        command=prev_card
    )
    prev_btn.pack(side="left")

    card_counter_label = tk.Label(nav_frame, text="Card 1 of 1", bg="white")
    card_counter_label.pack(side="left", padx=10)

    next_btn = tk.Button(
        nav_frame,
        text="Next ▶",
        font=("Arial", 9),
        bg="#999999",  # Changed to gray
        fg="white",
        padx=5,
        pady=1,
        relief="flat",
        command=next_card
    )
    next_btn.pack(side="left")

    # Configure the add card button
    add_card_btn.config(command=lambda: add_card_entry())

    # If editing, populate with existing cards
    if edit_deck and "cards" in edit_deck:
        for card in edit_deck["cards"]:
            add_card_entry(card.get("question", ""), card.get("answer", ""))
    else:
        # Add initial card entry
        add_card_entry()

    # Bottom buttons frame
    button_frame = tk.Frame(content_frame, bg="white")
    button_frame.pack(fill="x", pady=(10, 0))

    # Center the buttons
    button_center_frame = tk.Frame(button_frame, bg="white")
    button_center_frame.pack(side="top", fill="x")

    # Save Deck button
    save_btn = tk.Button(
        button_center_frame,
        text="Save Deck",
        font=("Arial", 10),
        bg="#4CAF50",  # Green color
        fg="white",
        padx=15,
        pady=5,
        relief="flat",
        command=lambda: save_manual_deck(title_entry, card_entries, card_window, edit_deck, edit_index)
    )
    save_btn.pack(side="left", padx=(0, 10))

    # Cancel button
    cancel_btn = tk.Button(
        button_center_frame,
        text="Cancel",
        font=("Arial", 10),
        bg="#999999",  # Gray color
        fg="white",
        padx=15,
        pady=5,
        relief="flat",
        command=card_window.destroy
    )
    cancel_btn.pack(side="left")


def save_manual_deck(title_entry, card_entries, card_window, edit_deck=None, edit_index=None):
    """Save a new deck or update an existing one"""
    global decks
    title = title_entry.get().strip()
    if not title:
        messagebox.showerror("Error", "Deck title is required.")
        return

    # Collect cards from entries
    manual_cards = []
    for entry in card_entries:
        question = entry["question"].get().strip()
        answer = entry["answer"].get().strip()

        if question and answer:  # Only add if both fields are filled
            manual_cards.append({
                "question": question,
                "answer": answer
            })

    if not manual_cards:
        messagebox.showerror("Error", "Please add at least one card with both question and answer.")
        return

    # Get current user
    username, _, _, _ = get_logged_in_user()

    if edit_deck and edit_index is not None:
        # Update existing deck
        decks[edit_index].update({
            "title": title,
            "cards": manual_cards,
            "card_count": len(manual_cards),
            "updated_at": time.strftime("%Y-%m-%d %H:%M:%S")
        })
        success_message = f"Deck '{title}' updated successfully with {len(manual_cards)} cards!"
        new_deck = decks[edit_index]
    else:
        # Create a new deck
        new_deck = {
            "title": title,
            "cards": manual_cards,
            "card_count": len(manual_cards),
            "created_by": username,
            "created_at": time.strftime("%Y-%m-%d %H:%M:%S"),
            "is_favorite": False  # Add favorite flag for filtering
        }
        # Add to decks list
        decks.append(new_deck)
        success_message = f"Deck '{title}' created successfully with {len(manual_cards)} cards!"

    # Save decks to file
    save_decks()

    # Update the dashboard
    update_deck_display()

    # Update progress display
    update_progress_display()

    messagebox.showinfo("Success", success_message)
    card_window.destroy()

    # Show the deck
    show_generated_flashcard_window(new_deck)


def show_generated_flashcard_window(deck=None):
    global flashcard_label, answer_label, current_index, flashcards, answer_hidden

    # Get user settings
    settings = get_current_settings()
    theme = settings["theme"]
    font_size = settings["font_size"]
    auto_reveal_time = settings["auto_reveal_time"]
    card_order = settings["card_order"]

    # If a deck is provided, use its flashcards
    if deck:
        flashcards = deck["cards"].copy()  # Make a copy to avoid modifying the original
        current_deck_title = deck["title"]

        # If random order is selected, shuffle the cards
        if card_order == "random":
            random.shuffle(flashcards)
    else:
        current_deck_title = "Unknown Deck"

    if not flashcards:
        messagebox.showerror("Error", "No flashcards available.")
        return

    # Initialize answer_hidden
    answer_hidden = True

    # Add this after initializing answer_hidden
    auto_reveal_timer_id = None

    # Add this function inside show_generated_flashcard_window
    def setup_auto_reveal_timer():
        nonlocal auto_reveal_timer_id
        # Cancel any existing timer
        if auto_reveal_timer_id is not None:
            flashcard_window.after_cancel(auto_reveal_timer_id)
            auto_reveal_timer_id = None

        # Set new timer if auto-reveal is enabled
        if auto_reveal_time > 0 and answer_hidden:
            auto_reveal_timer_id = flashcard_window.after(auto_reveal_time * 1000,
                                                          lambda: reveal_answer_with_feedback(feedback_frame))

    # Create the flashcard window
    flashcard_window = tk.Toplevel(root)
    flashcard_window.title("Flashcards")
    center_window(flashcard_window, 600, 500)

    # Create a dark header
    header_frame = tk.Frame(flashcard_window, bg="#333333", height=30)
    header_frame.pack(fill="x")

    # Add title to header
    header_title = tk.Label(header_frame, text="Flashcards", fg="white", bg="#333333", font=("Arial", 10))
    header_title.pack(side="left", padx=10, pady=5)

    # Add close button to header
    close_btn = tk.Button(
        header_frame,
        text="✕",
        font=("Arial", 10),
        bg="#333333",
        fg="white",
        bd=0,
        padx=10,
        command=flashcard_window.destroy
    )
    close_btn.pack(side="right")

    # Main container with padding
    main_container = tk.Frame(flashcard_window, bg="white", padx=20, pady=20)
    main_container.pack(fill="both", expand=True)

    # Title and card counter in the same row
    title_row = tk.Frame(main_container, bg="white")
    title_row.pack(fill="x", pady=(0, 10))

    # Title if deck is provided
    if deck:
        title_label = tk.Label(title_row, text=deck["title"], font=("Arial", 16, "bold"), bg="white")
        title_label.pack(side="left")

        info_label = tk.Label(title_row, text=f"Card {current_index + 1} of {len(flashcards)}",
                              font=("Arial", 12), bg="white", fg="#666666")
        info_label.pack(side="right")

    # Card content frame with border - takes most of the space
    card_frame = tk.Frame(main_container, bd=1, relief="solid", bg="#f2f2f2")
    card_frame.pack(fill="both", expand=True, pady=10)

    # Inner padding for card content
    card_content = tk.Frame(card_frame, padx=20, pady=20, bg="#f2f2f2")
    card_content.pack(fill="both", expand=True)

    # Question label - centered
    flashcard_label = tk.Label(card_content, text="", wraplength=500, justify="center",
                               font=("Arial", 14, "bold"), bg="#f2f2f2")
    flashcard_label.pack(fill="both", expand=True)

    # Answer label - centered
    answer_label = tk.Label(card_content, text="", wraplength=500, justify="center",
                            font=("Arial", 14), bg="#f2f2f2")
    answer_label.pack(fill="both", expand=True)

    # Feedback buttons frame - at the bottom of the card
    feedback_frame = tk.Frame(card_content, bg="#f2f2f2")
    feedback_frame.pack(fill="x", pady=(20, 0), side="bottom")

    # Make the buttons take equal width
    feedback_frame.columnconfigure(0, weight=1)
    feedback_frame.columnconfigure(1, weight=1)

    # Correct button - green with rounded corners
    correct_btn = tk.Button(
        feedback_frame,
        text="I Got It Right!",
        font=("Arial", 11),
        bg="#76c442",  # Bright green color
        fg="white",
        relief="flat",
        padx=15,
        pady=8,
        command=lambda: mark_answer(current_deck_title, True, info_label, feedback_frame)
    )
    correct_btn.grid(row=0, column=0, padx=(0, 5), sticky="ew")

    # Incorrect button - red with rounded corners
    incorrect_btn = tk.Button(
        feedback_frame,
        text="I Got It Wrong!",
        font=("Arial", 11),
        bg="#e74c3c",  # Bright red color
        fg="white",
        relief="flat",
        padx=15,
        pady=8,
        command=lambda: mark_answer(current_deck_title, False, info_label, feedback_frame)
    )
    incorrect_btn.grid(row=0, column=1, padx=(5, 0), sticky="ew")

    # Initially hide the feedback buttons
    feedback_frame.pack_forget()

    # Navigation buttons at the bottom
    nav_frame = tk.Frame(main_container, bg="white")
    nav_frame.pack(fill="x", pady=(10, 0))

    # Previous button (left arrow)
    prev_btn = tk.Button(
        nav_frame,
        text="<",
        font=("Arial", 16),
        bg="#eeeeee",
        relief="flat",
        width=3,
        command=lambda: prev_card_and_update_info(info_label, feedback_frame)
    )
    prev_btn.pack(side="left")

    # Next button (right arrow)
    next_btn = tk.Button(
        nav_frame,
        text=">",
        font=("Arial", 16),
        bg="#eeeeee",
        relief="flat",
        width=3,
        command=lambda: next_card_and_update_info(info_label, feedback_frame)
    )
    next_btn.pack(side="right")

    # Show Answer button - in the middle
    show_answer_btn = tk.Button(
        nav_frame,
        text="Show Answer",
        font=("Arial", 11),
        bg="#3498db",  # Blue color
        fg="white",
        relief="flat",
        padx=15,
        pady=5,
        command=lambda: reveal_answer_with_feedback(feedback_frame)
    )
    show_answer_btn.pack(side="top", pady=(5, 0))

    # Initialize the first card
    current_index = 0
    show_flashcard()
    setup_auto_reveal_timer()  # Set up timer for the first card

    # Function to show feedback buttons after revealing answer
    def reveal_answer_with_feedback(feedback_frame):
        reveal_answer()
        feedback_frame.pack(fill="x", pady=(20, 0), side="bottom")  # Show feedback buttons

    # Function to go to previous card
    def prev_card_and_update_info(info_label, feedback_frame):
        global current_index, answer_hidden
        if len(flashcards) > 1 and current_index > 0:
            current_index -= 1
            show_flashcard()
            answer_hidden = True
            info_label.config(text=f"Card {current_index + 1} of {len(flashcards)}")
            feedback_frame.pack_forget()  # Hide feedback buttons
            setup_auto_reveal_timer()  # Reset timer for the new card

    # Function to go to next card
    def next_card_and_update_info(info_label, feedback_frame):
        global current_index, answer_hidden
        if len(flashcards) > 1:
            current_index = (current_index + 1) % len(flashcards)
            show_flashcard()
            answer_hidden = True
            info_label.config(text=f"Card {current_index + 1} of {len(flashcards)}")
            feedback_frame.pack_forget()  # Hide feedback buttons
            setup_auto_reveal_timer()  # Reset timer for the new card


def mark_answer(deck_title, is_correct, info_label, feedback_frame):
    """Mark the current card as answered correctly or incorrectly"""
    update_card_progress(deck_title, current_index, is_correct)
    next_card_and_update_info(info_label, feedback_frame)


def next_card_and_update_info(info_label=None, feedback_frame=None):
    global current_index, answer_hidden
    if flashcards:
        current_index = (current_index + 1) % len(flashcards)
        show_flashcard()
        answer_hidden = True  # Reset answer hidden state
        if info_label:
            info_label.config(text=f"Card {current_index + 1} of {len(flashcards)}")
        if feedback_frame:
            feedback_frame.pack_forget()  # Hide feedback buttons


def toggle_favorite(deck_index):
    """Toggle the favorite status of a deck"""
    if 0 <= deck_index < len(decks):
        # Toggle the favorite status
        decks[deck_index]['is_favorite'] = not decks[deck_index].get('is_favorite', False)
        # Save the updated decks list
        save_decks()
        # Update the display with current filter
        update_deck_display(current_filter.get())


def create_deck_card(parent, deck, index):
    """Create a visual card for a deck in the dashboard"""
    # Get user settings
    settings = get_current_settings()
    theme = settings["theme"]
    font_size = settings["font_size"]

    # Create card frame with white background and border - now more rectangular and wider
    card_frame = tk.Frame(parent, bd=1, relief=tk.SOLID, bg="white", width=700, height=80)
    card_frame.pack_propagate(False)  # Keep fixed size

    # Left side - Title and info
    info_frame = tk.Frame(card_frame, bg="white", padx=15, pady=10)
    info_frame.pack(side="left", fill="both", expand=True)

    # Title and card count in the same row
    title_row = tk.Frame(info_frame, bg="white")
    title_row.pack(fill="x", anchor="w")

    # Title
    title_label = tk.Label(title_row, text=deck["title"], font=("Arial", 12, "bold"), bg="white")
    title_label.pack(side="left", anchor="w")

    # Card count
    count_label = tk.Label(title_row, text=f"{deck['card_count']} Cards", font=("Arial", 10),
                           bg="white", fg="#FFD700")
    count_label.pack(side="left", padx=(10, 0), anchor="w")

    # Favorite star (if favorite)
    if deck.get('is_favorite', False):
        fav_label = tk.Label(title_row, text="★", font=("Arial", 14), bg="white", fg="#FFD700")
        fav_label.pack(side="left", padx=(5, 0), anchor="w")

    # Created date in a cleaner format
    if "created_at" in deck:
        # Parse the date string and format it as "MM-DD-YYYY"
        try:
            date_obj = time.strptime(deck['created_at'], "%Y-%m-%d %H:%M:%S")
            formatted_date = time.strftime("%m-%d-%Y", date_obj)
        except:
            formatted_date = deck['created_at']

        date_label = tk.Label(info_frame, text=f"Created: {formatted_date}",
                              font=("Arial", 9), bg="white", fg="#666666")
        date_label.pack(anchor="w", pady=(5, 0))

    # Right side - Buttons
    buttons_frame = tk.Frame(card_frame, bg="white", padx=15, pady=10)
    buttons_frame.pack(side="right")

    # Favorite button
    fav_text = "★ Unfavorite" if deck.get('is_favorite', False) else "☆ Favorite"
    fav_btn = tk.Button(
        buttons_frame,
        text=fav_text,
        font=("Arial", 10),
        bg="#FFD700" if deck.get('is_favorite', False) else "#EEEEEE",
        fg="white" if deck.get('is_favorite', False) else "black",
        relief="flat",
        padx=10,
        pady=2,
        command=lambda idx=index: toggle_favorite(idx)
    )
    fav_btn.pack(side=tk.RIGHT, padx=(5, 0))

    # Edit button - blue
    edit_btn = tk.Button(
        buttons_frame,
        text="Edit",
        font=("Arial", 10),
        bg="#2196F3",
        fg="white",
        relief="flat",
        padx=10,
        pady=2,
        command=lambda d=deck, idx=index: edit_deck(d, idx)
    )
    edit_btn.pack(side=tk.RIGHT, padx=(5, 0))

    # Delete button - red
    delete_btn = tk.Button(
        buttons_frame,
        text="Delete",
        font=("Arial", 10),
        bg="#F44336",
        fg="white",
        relief="flat",
        padx=10,
        pady=2,
        command=lambda idx=index: delete_deck(idx)
    )
    delete_btn.pack(side=tk.RIGHT)

    # Make the entire card clickable to open the deck
    card_frame.bind("<Button-1>", lambda e, d=deck: show_generated_flashcard_window(d))
    for widget in info_frame.winfo_children():
        widget.bind("<Button-1>", lambda e, d=deck: show_generated_flashcard_window(d))
    for child in info_frame.winfo_children():
        if isinstance(child, tk.Frame):
            for widget in child.winfo_children():
                widget.bind("<Button-1>", lambda e, d=deck: show_generated_flashcard_window(d))

    return card_frame


def update_deck_display(filter_option="all"):
    """Update the deck display in the dashboard with filtering"""
    # Clear existing deck display
    for widget in decks_container.winfo_children():
        widget.destroy()

    # Get current user
    username, _, _, _ = get_logged_in_user()

    # Filter decks to only show those created by the current user
    user_decks = [deck for deck in decks if deck.get("created_by") == username]

    # Apply filter
    if filter_option == "favorites":
        user_decks = [deck for deck in user_decks if deck.get("is_favorite", False)]
    elif filter_option == "recent":
        # Sort by creation date (newest first)
        user_decks = sorted(user_decks, key=lambda d: d.get("created_at", ""), reverse=True)
        # Limit to 5 most recent
        user_decks = user_decks[:5]

    # Get user settings
    settings = get_current_settings()
    theme = settings["theme"]
    font_size = settings["font_size"]

    # If no decks, show a message
    if not user_decks:
        no_decks_label = tk.Label(decks_container, text="No flashcard decks yet. Create one with the Add Card button!")
        apply_theme(no_decks_label, theme)
        apply_font_size(no_decks_label, font_size, "regular")
        no_decks_label.pack(pady=20)
        return

    # Create deck cards in a vertical list
    for i, deck in enumerate(user_decks):
        card = create_deck_card(decks_container, deck, decks.index(deck))  # Use the index in the full decks list
        card.pack(fill="x", pady=5, padx=10)


def edit_deck(deck, index):
    """Edit an existing deck"""
    open_card_window(deck, index)


def update_progress_display():
    """Update the progress display in the dashboard"""
    # Calculate progress percentage based on decks and correct answers
    username = get_logged_in_user()[0]
    progress_data = get_user_progress(username)

    # Count total cards, right answers, and wrong answers
    total_cards = 0
    right_cards = 0
    wrong_cards = 0

    for deck in decks:
        if deck.get("created_by") == username:  # Only count user's decks
            deck_title = deck.get('title', '')
            card_count = deck.get('card_count', 0)
            total_cards += card_count

            # Get progress for this deck
            deck_progress = progress_data.get(deck_title, {})
            card_status = deck_progress.get('card_status', {})

            # Count cards based on their current status
            deck_right = 0
            deck_wrong = 0

            for status in card_status.values():
                if status == 'correct':
                    deck_right += 1
                elif status == 'wrong':
                    deck_wrong += 1

            right_cards += deck_right
            wrong_cards += deck_wrong

    # Calculate progress percentage (limit to 100%)
    progress_percentage = min(100, (right_cards / max(total_cards, 1) * 100))

    # Update the progress display
    if 'progress_canvas' in globals():
        # Clear the canvas
        progress_canvas.delete("all")

        # Draw background circle (light gray)
        progress_canvas.create_oval(10, 10, 90, 90, outline="#E0E0E0", width=8, fill="white")

        # Draw progress arc (gold color to match screenshot)
        start_angle = 90  # Start from top (90 degrees)
        extent_angle = 3.6 * progress_percentage  # 3.6 degrees per percentage point (360 / 100)
        progress_canvas.create_arc(10, 10, 90, 90, start=start_angle, extent=extent_angle,
                                   outline="#FFD700", width=8, style="arc")

        # Add percentage text in the center - rounded to whole number
        progress_canvas.create_text(50, 50, text=f"{int(progress_percentage)}%", font=("Arial", 12, "bold"))

    # Update the stats
    if 'total_number' in globals():
        total_number.config(text=str(total_cards))

    if 'right_number' in globals():
        right_number.config(text=str(right_cards))

    if 'wrong_number' in globals():
        wrong_number.config(text=str(wrong_cards))


def open_settings():
    """Open the settings page"""
    open_page("settings")


# Main application setup
if __name__ == "__main__":
    root = tk.Tk()
    root.title("FlashLearn - Dashboard")

    # Set the window size
    window_width = 1000
    window_height = 700

    # Center the main window
    center_window(root, window_width, window_height)

    # Get user data
    logged_in_username, logged_in_email, logged_in_password, profile_picture_path = get_logged_in_user()

    # Get user settings
    settings = get_current_settings()
    theme = settings["theme"]
    font_size = settings["font_size"]

    # Apply theme to root window
    apply_theme(root, theme)
    root.resizable(False, False)

    # Variable to track current filter
    current_filter = tk.StringVar(value="all")


    # Function to load and display the profile icon
    def display_profile_icon():
        if profile_picture_path and os.path.exists(profile_picture_path):
            try:
                img = Image.open(profile_picture_path)
                img = img.resize((50, 50), Image.LANCZOS)

                # Make the image circular
                mask = Image.new('L', (50, 50), 0)
                draw = ImageDraw.Draw(mask)
                draw.ellipse((0, 0, 50, 50), fill=255)
                img_circular = Image.new('RGBA', (50, 50), (0, 0, 0, 0))
                img_circular.paste(img, (0, 0), mask)

                img_tk = ImageTk.PhotoImage(img_circular)
                profile_icon_label.config(image=img_tk)
                profile_icon_label.image = img_tk
            except Exception as e:
                print(f"Error displaying profile icon: {e}")
                default_image = Image.new("RGB", (50, 50), color="#CCCCCC")
                default_image = ImageTk.PhotoImage(default_image)
                profile_icon_label.config(image=default_image)
                profile_icon_label.image = default_image
        else:
            default_image = Image.new("RGB", (50, 50), color="#CCCCCC")
            default_image = ImageTk.PhotoImage(default_image)
            profile_icon_label.config(image=default_image)
            profile_icon_label.image = default_image


    # Sidebar Frame (Fixed Width)
    sidebar = tk.Frame(root, width=220, height=600, highlightthickness=0 if theme == "dark" else 2)
    apply_theme(sidebar, theme, "sidebar")
    sidebar.pack_propagate(False)  # Prevent resizing
    sidebar.pack(side="left", fill="y")

    # FlashLearn Title
    title_label = tk.Label(sidebar, text="FlashLearn")
    apply_theme(title_label, theme, "sidebar")
    apply_font_size(title_label, font_size, "title")
    title_label.pack(pady=20)

    # Sidebar Buttons Configuration
    menu_items = {
        "⚫ Profile": "profile",
        "📖 My Cards": None,
        "🔥 Progress": "progress",
        "⚙️ Settings": "settings"  # Now points to settings.py
    }

    buttons = {}


    def select_menu(selected):
        """ Highlights the selected menu item and resets others """
        for btn_text, (btn, page) in buttons.items():
            if btn_text == selected:
                apply_theme(btn, theme, "highlight")
                apply_font_size(btn, font_size, "subtitle")
            else:
                apply_theme(btn, theme, "sidebar")
                apply_font_size(btn, font_size, "regular")


    # Create Sidebar Buttons with Navigation
    for text, page in menu_items.items():
        btn = tk.Button(
            sidebar, text=text, bd=0,
            anchor="w", padx=20, width=20,
            command=lambda t=text, p=page: open_page(p) if p else select_menu(t)
        )
        apply_theme(btn, theme, "sidebar")
        apply_font_size(btn, font_size, "regular")
        btn.pack(pady=5)
        buttons[text] = (btn, page)  # Store button references

    # Default: Select "My Cards"
    select_menu("📖 My Cards")

    # Sign Out Button (Goes back to sign_in.py)
    sign_out_btn = tk.Button(sidebar, text="↩ Sign Out", bd=0,
                             anchor="w", padx=20, width=20, command=sign_out)
    apply_theme(sign_out_btn, theme, "sidebar")
    apply_font_size(sign_out_btn, font_size, "regular")
    sign_out_btn.pack(pady=30, side="bottom")

    # Main Content Frame
    content_frame = tk.Frame(root, width=680, height=600)
    apply_theme(content_frame, theme)
    content_frame.pack(side="right", fill="both", expand=True)

    # Dashboard Title
    dashboard_label = tk.Label(content_frame, text="Dashboard")
    apply_theme(dashboard_label, theme)
    apply_font_size(dashboard_label, font_size, "title")
    dashboard_label.pack(pady=20)

    # Profile Icon Section (Upper Right)
    profile_icon_label = tk.Label(root)
    apply_theme(profile_icon_label, theme)
    profile_icon_label.place(x=window_width - 60, y=10)  # Place icon 10px from the top-right corner

    # Call the function to display the profile icon
    display_profile_icon()

    # Dashboard Title and Profile Icon in the same Frame
    title_and_icon_frame = tk.Frame(content_frame)
    apply_theme(title_and_icon_frame, theme)
    title_and_icon_frame.pack(fill="x", padx=20, pady=10)

    # Greeting Section
    greeting_frame = tk.Frame(content_frame, padx=20, pady=20)
    apply_theme(greeting_frame, theme, "card")
    greeting_frame.pack(pady=10, padx=20, fill="x")

    # Update the greeting label to reflect the logged-in username
    greeting_label = tk.Label(greeting_frame, text=f"Welcome Back, {logged_in_username}!",
                              font=("Arial", 14, "bold"), bg="white")
    apply_theme(greeting_label, theme, "card")
    greeting_label.pack(anchor="w")

    # Description text with wrapping and justification
    description_text = """Get ready to boost your knowledge! Here, you can create, review, and manage your flashcards to help
    you learn and memorize key concepts effectively. You're in the right place!"""

    # Wrap the text to fit within the given width (600px)
    wrapped_text = textwrap.fill(description_text, width=95)  # Adjust width for your desired wrapping

    # Description label with the regular text below the greeting
    desc_label = tk.Label(greeting_frame, text=wrapped_text, justify="left", wraplength=600)
    apply_theme(desc_label, theme, "card")
    apply_font_size(desc_label, font_size, "regular")
    desc_label.pack(anchor="w", padx=10, pady=5)

    # Add Card Button
    add_card_btn = tk.Button(greeting_frame, text="➕ Add Card", padx=10, pady=5, command=open_card_window)
    apply_theme(add_card_btn, theme, "button")
    apply_font_size(add_card_btn, font_size, "button")
    add_card_btn.pack(anchor="w", pady=5)

    # Progress Section (New)
    progress_frame = tk.Frame(content_frame, padx=20, pady=20, bg="white", bd=1, relief="solid")
    progress_frame.pack(pady=10, padx=20, fill="x")

    # Progress title
    progress_title = tk.Label(progress_frame, text="Progress", font=("Arial", 14, "bold"), bg="white")
    progress_title.pack(anchor="w", pady=(0, 10))

    # Progress content frame (holds circular progress and stats)
    progress_content = tk.Frame(progress_frame, bg="white")
    progress_content.pack(fill="x")

    # Left side - Circular progress indicator
    progress_circle_frame = tk.Frame(progress_content, bg="white", width=100, height=100)
    progress_circle_frame.pack(side="left", padx=20)
    progress_circle_frame.pack_propagate(False)  # Keep fixed size

    # Create canvas for circular progress
    progress_canvas = tk.Canvas(progress_circle_frame, width=100, height=100, bg="white", highlightthickness=0)
    progress_canvas.pack()

    # Right side - Stats
    stats_frame = tk.Frame(progress_content, bg="white")
    stats_frame.pack(side="left", padx=20, fill="x", expand=True)

    # Total Cards
    total_frame = tk.Frame(stats_frame, bg="white")
    total_frame.pack(side="left", padx=20)

    total_number = tk.Label(total_frame, text="0", font=("Arial", 18, "bold"), bg="white")
    total_number.pack()

    total_label = tk.Label(total_frame, text="Total Cards", font=("Arial", 10), bg="white")
    total_label.pack()

    # Right Answers
    right_frame = tk.Frame(stats_frame, bg="white")
    right_frame.pack(side="left", padx=20)

    right_number = tk.Label(right_frame, text="0", font=("Arial", 18, "bold"), bg="white", fg="#4CAF50")
    right_number.pack()

    right_label = tk.Label(right_frame, text="Right", font=("Arial", 10), bg="white")
    right_label.pack()

    # Wrong Answers
    wrong_frame = tk.Frame(stats_frame, bg="white")
    wrong_frame.pack(side="left", padx=20)

    wrong_number = tk.Label(wrong_frame, text="0", font=("Arial", 18, "bold"), bg="white", fg="#F44336")
    wrong_number.pack()

    wrong_label = tk.Label(wrong_frame, text="Wrong", font=("Arial", 10), bg="white")
    wrong_label.pack()

    # My Decks Section with Filter dropdown
    decks_header_frame = tk.Frame(content_frame)
    apply_theme(decks_header_frame, theme)
    decks_header_frame.pack(fill="x", padx=20, pady=5)

    my_decks_label = tk.Label(decks_header_frame, text="My Decks", font=("Arial", 14, "bold"))
    apply_theme(my_decks_label, theme)
    my_decks_label.pack(side="left")


    # Filter dropdown menu
    def apply_filter(option):
        current_filter.set(option)
        update_deck_display(option)


    filter_frame = tk.Frame(decks_header_frame)
    filter_frame.pack(side="right")

    filter_label = tk.Label(filter_frame, text="Filter: ")
    filter_label.pack(side="left")

    filter_options = ["all", "recent", "favorites"]
    filter_menu = ttk.Combobox(filter_frame, textvariable=current_filter, values=filter_options,
                               state="readonly", width=10)
    filter_menu.pack(side="left")
    filter_menu.bind("<<ComboboxSelected>>", lambda e: apply_filter(current_filter.get()))

    # Create a scrollable frame for decks
    decks_frame = tk.Frame(content_frame)
    apply_theme(decks_frame, theme)
    decks_frame.pack(fill="both", expand=True, padx=20, pady=5)

    # Canvas for scrolling
    canvas = tk.Canvas(decks_frame)
    apply_theme(canvas, theme)
    scrollbar = tk.Scrollbar(decks_frame, orient="vertical", command=canvas.yview)
    canvas.configure(yscrollcommand=scrollbar.set)

    # Pack scrollbar and canvas
    scrollbar.pack(side="right", fill="y")
    canvas.pack(side="left", fill="both", expand=True)

    # Create a frame inside the canvas to hold the deck cards
    decks_container = tk.Frame(canvas)
    apply_theme(decks_container, theme)
    canvas.create_window((0, 0), window=decks_container, anchor="nw")


    # Configure the canvas to resize with the window
    def configure_canvas(event):
        canvas.configure(scrollregion=canvas.bbox("all"), width=event.width)
        decks_container.configure(width=event.width)


    canvas.bind("<Configure>", configure_canvas)
    decks_container.bind("<Configure>", lambda e: canvas.configure(scrollregion=canvas.bbox("all")))

    # Load existing decks
    load_decks()
    update_progress_display()

    # Update the deck display
    update_deck_display()

    # Update the progress display
    update_progress_display()

    # Run the Application
    root.mainloop()
